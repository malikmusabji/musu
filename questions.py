import streamlit as st
from utils import get_llminfo
from pypdf import PdfReader
from pptx import Presentation
import google.generativeai as genai

def questions_page():
    st.subheader("Questions Page")
    st.markdown("""The Questions page allows you to upload PDFs and PPTs, automatically extracting the information and generating interactive MCQ questions for practice. To your left is parameter control for the LLM you chose to use.""")
    model, temperature, top_p, max_tokens, top_k = get_llminfo()

    uploaded_file = st.file_uploader("Upload a PDF or PPT file", type=["pdf", "ppt", "pptx"])

    if uploaded_file is not None:
        text = ""
        if uploaded_file.type == "application/pdf":
            pdf_reader = PdfReader(uploaded_file)
            for page in pdf_reader.pages:
                text += page.extract_text()
        elif uploaded_file.type in ["application/vnd.ms-powerpoint", "application/vnd.openxmlformats-officedocument.presentationml.presentation"]:
            presentation = Presentation(uploaded_file)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            text = text.strip()

        if st.button("Generate MCQs"):
            if text:
                generation_config = {
                    "temperature": temperature,
                    "top_p": top_p,
                    "max_output_tokens": max_tokens,
                    "top_k": top_k
                }
                model_instance = genai.GenerativeModel(model_name=model, generation_config=generation_config)

                prompt = (
                    f"Please generate exactly 5 multiple-choice questions based on the following text:\n\n"
                    f"{text}\n\n"
                    "Each question should have 4 options labeled A, B, C, and D. "
                    "At the end of each question, specify the correct answer in the format:\n"
                    "Correct Answer: [A/B/C/D]\n\n"
                    "For example:\n"
                    "1. What is the capital of France?\n"
                    "A) Paris\n"
                    "B) London\n"
                    "C) Rome\n"
                    "D) Berlin\n"
                    "Correct Answer: A\n\n"
                    "Now generate 5 questions following this format:"
                )
                
                response = model_instance.generate_content([prompt])
                mcqs_with_answers = response.text.strip().split('\n\n')

                st.session_state.mcqs = []
                st.session_state.correct_answers = []

                for mcq in mcqs_with_answers:
                    lines = mcq.split('\n')
                    if len(lines) < 2:
                        continue
                    
                    question_text = lines[0].strip()
                    options = [option.strip() for option in lines[1:5] if option.strip()]

                    correct_answer_line = lines[-1] if len(lines) > 5 else ""
                    correct_answer = correct_answer_line.split(":")[-1].strip().upper()
                    correct_answer = correct_answer.replace("**", "")

                    if not options or not correct_answer:
                        continue
                    
                    st.session_state.mcqs.append((question_text, options))
                    st.session_state.correct_answers.append(correct_answer)

                st.session_state.user_answers = [None] * len(st.session_state.mcqs)

        if 'mcqs' in st.session_state:
            st.subheader("Generated MCQs:")
            
            for i, (question_text, options) in enumerate(st.session_state.mcqs):
                selected_option = st.radio(
                    question_text, 
                    options, 
                    key=f"question_{i}", 
                    index=options.index(st.session_state.user_answers[i]) if st.session_state.user_answers[i] in options else 0
                )
                st.session_state.user_answers[i] = selected_option

            if st.button("Submit Answers"):
                correct_answers_count = 0

                for i, selected_option in enumerate(st.session_state.user_answers):
                    normalized_selected_option = selected_option[0].upper()
                    normalized_correct_answer = st.session_state.correct_answers[i].upper()

                    if normalized_selected_option == normalized_correct_answer:
                        correct_answers_count += 1
                
                total_questions = len(st.session_state.mcqs)
                st.success(f"You got {correct_answers_count} out of {total_questions} correct!")
