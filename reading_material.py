import streamlit as st
from utils import get_llminfo
from pypdf import PdfReader
from pptx import Presentation
import google.generativeai as genai
import time

def reading_material_page():
    st.subheader("Reading Material Interaction")
    st.markdown("""The Reading Material page allows you to upload various types of media, enabling you to chat with a chatbot about the content for better understanding and clarity also providing additional sources to read. To your left is parameter control for the LLM you chose to use.""")
    model, temperature, top_p, max_tokens, top_k = get_llminfo()

    typepdf = st.radio("Select the type of media to interact with:", ("PDF", "Images", "Videos", "PPT"), index=0)

    if typepdf == "PDF":
        st.write("You selected PDF. Upload your files below.")
        uploaded_files = st.file_uploader("Choose one or more PDFs", type='pdf', accept_multiple_files=True)
        if uploaded_files:
            text = ""
            for pdf in uploaded_files:
                pdf_reader = PdfReader(pdf)
                for page in pdf_reader.pages:
                    text += page.extract_text()

            generation_config = {
                "temperature": temperature,
                "top_p": top_p,
                "max_output_tokens": max_tokens,
                "top_k": top_k,
                "response_mime_type": "text/plain",
            }
            model_instance = genai.GenerativeModel(
                model_name=model,
                generation_config=generation_config,
            )
            st.write(model_instance.count_tokens(text))
            question = st.text_input("Enter your question and hit return.")
            if question:
                response = model_instance.generate_content([question, text])
                st.write(response.text)

    elif typepdf == "Images":
        st.write("You selected Images. Upload your image file below.")
        image_file = st.file_uploader("Upload your image file.", type=["jpg", "jpeg", "png"])
        if image_file:
            temp_file_path = image_file.name
            with open(temp_file_path, "wb") as f:
                f.write(image_file.getbuffer())

            st.write("Uploading image...")
            uploaded_image = genai.upload_file(path=temp_file_path)
            while uploaded_image.state.name == "PROCESSING":
                time.sleep(5)
                uploaded_image = genai.get_file(uploaded_image.name)

            if uploaded_image.state.name == "FAILED":
                st.error("Failed to process image.")
                return

            st.write("Image uploaded successfully. Enter your prompt below.")
            prompt2 = st.text_input("Enter your prompt for the image.")
            if prompt2:
                generation_config = {
                    "temperature": temperature,
                    "top_p": top_p,
                    "max_output_tokens": max_tokens,
                    "top_k": top_k,
                }
                response_image = model_instance.generate_content([prompt2, uploaded_image])
                st.write(response_image.text)

    elif typepdf == "Videos":
        st.write("You selected Videos. Upload your video file below.")
        video_file = st.file_uploader("Upload your video file.", type=["mp4", "mov", "avi"])
        if video_file:
            temp_file_path = video_file.name
            with open(temp_file_path, "wb") as f:
                f.write(video_file.getbuffer())

            st.write("Uploading video...")
            uploaded_video = genai.upload_file(path=temp_file_path)
            while uploaded_video.state.name == "PROCESSING":
                time.sleep(5)
                uploaded_video = genai.get_file(uploaded_video.name)

            if uploaded_video.state.name == "FAILED":
                st.error("Failed to process video.")
                return
            
            st.write("Video uploaded successfully. Enter your prompt below.")
            prompt3 = st.text_input("Enter your prompt for the video.")
            if prompt3:
                model_instance = genai.GenerativeModel(model_name=model)
                response = model_instance.generate_content([uploaded_video, prompt3])
                st.markdown(response.text)
                genai.delete_file(uploaded_video.name)

    elif typepdf == "PPT":
        st.write("You selected PPT. Upload your PowerPoint file below.")
        uploaded_ppt = st.file_uploader("Choose a PPT file", type='pptx')
        if uploaded_ppt:
            text = ""
            presentation = Presentation(uploaded_ppt)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            text = text.strip()
            st.write("Extracted text from PowerPoint:")
            st.write(text)

            generation_config = {
                "temperature": temperature,
                "top_p": top_p,
                "max_output_tokens": max_tokens,
                "top_k": top_k,
                "response_mime_type": "text/plain",
            }
            model_instance = genai.GenerativeModel(
                model_name=model,
                generation_config=generation_config,
            )
            st.write(model_instance.count_tokens(text))
            question = st.text_input("Enter your question about the PPT content and hit return.")
            if question:
                response = model_instance.generate_content([question, text])
                st.write(response.text)
